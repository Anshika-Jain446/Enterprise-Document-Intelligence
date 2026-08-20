import io
import os
from pathlib import Path

import pandas as pd
import streamlit as st

from werkzeug.security import (
    generate_password_hash,
    check_password_hash,
)

from chunking import ChunkingEngine
from model import GeminiModel


# ============================================================
# CONFIGURATION
# ============================================================

try:
    import config
except ImportError:
    config = None


def config_value(name, default=None):
    value = os.getenv(name)

    if value is not None:
        return value

    if config is not None:
        return getattr(config, name, default)

    return default


OPENROUTER_API_KEY = str(
    config_value("OPENROUTER_API_KEY", "") or ""
).strip()

LLM_MODEL = str(
    config_value(
        "LLM_MODEL",
        "google/gemini-2.0-flash-001",
    )
    or "google/gemini-2.0-flash-001"
).strip()

POSTGRES_HOST = str(
    config_value("POSTGRES_HOST", "localhost")
    or "localhost"
).strip()

try:
    POSTGRES_PORT = int(
        config_value("POSTGRES_PORT", 5432) or 5432
    )
except Exception:
    POSTGRES_PORT = 5432

POSTGRES_DB = str(
    config_value("POSTGRES_DB", "") or ""
).strip()

POSTGRES_USER = str(
    config_value("POSTGRES_USER", "") or ""
).strip()

POSTGRES_PASSWORD = str(
    config_value("POSTGRES_PASSWORD", "") or ""
)

try:
    CHUNK_SIZE = int(
        config_value("CHUNK_SIZE", 1000) or 1000
    )
except Exception:
    CHUNK_SIZE = 1000

try:
    CHUNK_OVERLAP = int(
        config_value("CHUNK_OVERLAP", 100) or 100
    )
except Exception:
    CHUNK_OVERLAP = 100


SUPPORTED_FILES_CONFIG = config_value(
    "SUPPORTED_FILES",
    [
        ".pdf",
        ".docx",
        ".pptx",
        ".xlsx",
        ".xls",
        ".csv",
        ".txt",
        ".md",
        ".markdown",
    ],
)

if isinstance(SUPPORTED_FILES_CONFIG, str):
    SUPPORTED_FILES = [
        item.strip()
        for item in SUPPORTED_FILES_CONFIG.split(",")
        if item.strip()
    ]
else:
    SUPPORTED_FILES = list(
        SUPPORTED_FILES_CONFIG or []
    )


SUPPORTED_EXTENSIONS = {
    (
        str(item).lower()
        if str(item).startswith(".")
        else f".{str(item).lower()}"
    )
    for item in SUPPORTED_FILES
}


DATABASE_URL = str(
    os.getenv(
        "DATABASE_URL",
        os.getenv("POSTGRESQL_URL", ""),
    )
    or ""
).strip()

ADMIN_USERNAME = str(
    os.getenv("ADMIN_USERNAME", "") or ""
).strip()

DEFAULT_CHUNK_SIZE = CHUNK_SIZE
DEFAULT_CHUNK_OVERLAP = CHUNK_OVERLAP
DEFAULT_TOP_K = 5


# ============================================================
# OPTIONAL POSTGRESQL IMPORTS
# ============================================================

try:
    import psycopg2
    from psycopg2.extras import RealDictCursor, Json
except ImportError:
    psycopg2 = None
    RealDictCursor = None
    Json = None


try:
    from pgvector.psycopg2 import register_vector
except ImportError:
    register_vector = None


# ============================================================
# DOCUMENT READER IMPORTS
# ============================================================

try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None


try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None


try:
    from pptx import Presentation
except ImportError:
    Presentation = None


# ============================================================
# STREAMLIT CONFIG
# ============================================================

st.set_page_config(
    page_title="Enterprise Document Intelligence",
    page_icon="📄",
    layout="wide",
    initial_sidebar_state="expanded",
)


# ============================================================
# DATABASE STORE
# ============================================================

class PostgreSQLStore:

    def __init__(self):
        self.connection = None

    # ========================================================
    # CONNECTION
    # ========================================================

    def connect(self):

        if psycopg2 is None:
            raise RuntimeError(
                "psycopg2 is not installed. "
                "Install psycopg2-binary."
            )

        # Reuse healthy connection.
        if self.connection is not None:

            try:
                if self.connection.closed == 0:
                    return self.connection
            except Exception:
                pass

            self.reset_connection()

        try:

            if DATABASE_URL:

                self.connection = psycopg2.connect(
                    DATABASE_URL,
                    connect_timeout=15,
                )

            else:

                if not POSTGRES_DB:
                    raise RuntimeError(
                        "PostgreSQL is not configured. "
                        "Set DATABASE_URL or POSTGRES_DB."
                    )

                self.connection = psycopg2.connect(
                    host=POSTGRES_HOST,
                    port=POSTGRES_PORT,
                    database=POSTGRES_DB,
                    user=POSTGRES_USER,
                    password=POSTGRES_PASSWORD,
                    connect_timeout=15,
                )

            self.connection.autocommit = False

            if register_vector is not None:

                try:
                    register_vector(self.connection)
                except Exception:
                    pass

            cursor = self.connection.cursor()

            try:

                cursor.execute(
                    "SET search_path TO public"
                )

                self.connection.commit()

            except Exception:

                self.connection.rollback()
                raise

            finally:
                cursor.close()

            return self.connection

        except Exception:

            self.reset_connection()
            raise

    # ========================================================
    # RESET CONNECTION
    # ========================================================

    def reset_connection(self):

        try:

            if self.connection is not None:
                self.connection.close()

        except Exception:
            pass

        self.connection = None

    # ========================================================
    # INITIALIZE DATABASE
    # ========================================================

    def initialize(self):

        conn = self.connect()
        cursor = None

        try:

            cursor = conn.cursor()

            # ------------------------------------------------
            # USERS
            # ------------------------------------------------

            cursor.execute(
                """
                CREATE TABLE IF NOT EXISTS public.users (
                    id SERIAL PRIMARY KEY,
                    username VARCHAR(255) UNIQUE NOT NULL,
                    password_hash TEXT NOT NULL,
                    role VARCHAR(50) NOT NULL DEFAULT 'user',
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
                """
            )

            # ------------------------------------------------
            # DOCUMENTS
            # ------------------------------------------------

            cursor.execute(
                """
                CREATE TABLE IF NOT EXISTS public.documents (
                    id SERIAL PRIMARY KEY,
                    user_id INTEGER NOT NULL,
                    filename TEXT NOT NULL,
                    file_type TEXT,
                    file_size BIGINT,
                    file_data BYTEA,
                    chunking_method TEXT,
                    chunk_size INTEGER,
                    chunk_overlap INTEGER,
                    metadata JSONB DEFAULT '{}'::jsonb,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,

                    CONSTRAINT documents_user_fk
                    FOREIGN KEY (user_id)
                    REFERENCES public.users(id)
                    ON DELETE CASCADE
                )
                """
            )

            # ------------------------------------------------
            # CHUNKS
            # ------------------------------------------------

            cursor.execute(
                """
                CREATE TABLE IF NOT EXISTS public.document_chunks (
                    id SERIAL PRIMARY KEY,
                    document_id INTEGER NOT NULL,
                    chunk_id INTEGER NOT NULL,
                    chunk_type TEXT,
                    content TEXT,
                    page TEXT,
                    tokens INTEGER,
                    characters INTEGER,
                    metadata JSONB DEFAULT '{}'::jsonb,
                    chunk_data JSONB DEFAULT '{}'::jsonb,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,

                    CONSTRAINT chunks_document_fk
                    FOREIGN KEY (document_id)
                    REFERENCES public.documents(id)
                    ON DELETE CASCADE,

                    CONSTRAINT document_chunk_unique
                    UNIQUE(document_id, chunk_id)
                )
                """
            )

            # ------------------------------------------------
            # CONVERSATIONS
            # ------------------------------------------------

            cursor.execute(
                """
                CREATE TABLE IF NOT EXISTS public.conversations (
                    id SERIAL PRIMARY KEY,
                    user_id INTEGER NOT NULL,
                    title TEXT DEFAULT 'New Conversation',
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,

                    CONSTRAINT conversations_user_fk
                    FOREIGN KEY (user_id)
                    REFERENCES public.users(id)
                    ON DELETE CASCADE
                )
                """
            )

            # ------------------------------------------------
            # MESSAGES
            # ------------------------------------------------

            cursor.execute(
                """
                CREATE TABLE IF NOT EXISTS public.messages (
                    id SERIAL PRIMARY KEY,
                    conversation_id INTEGER NOT NULL,
                    role TEXT NOT NULL,
                    content TEXT NOT NULL,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,

                    CONSTRAINT messages_conversation_fk
                    FOREIGN KEY (conversation_id)
                    REFERENCES public.conversations(id)
                    ON DELETE CASCADE
                )
                """
            )

            # ------------------------------------------------
            # MIGRATIONS
            # ------------------------------------------------

            cursor.execute(
                """
                ALTER TABLE public.documents
                ADD COLUMN IF NOT EXISTS metadata
                JSONB DEFAULT '{}'::jsonb
                """
            )

            cursor.execute(
                """
                ALTER TABLE public.documents
                ADD COLUMN IF NOT EXISTS chunking_method TEXT
                """
            )

            cursor.execute(
                """
                ALTER TABLE public.documents
                ADD COLUMN IF NOT EXISTS chunk_size INTEGER
                """
            )

            cursor.execute(
                """
                ALTER TABLE public.documents
                ADD COLUMN IF NOT EXISTS chunk_overlap INTEGER
                """
            )

            cursor.execute(
                """
                ALTER TABLE public.document_chunks
                ADD COLUMN IF NOT EXISTS metadata
                JSONB DEFAULT '{}'::jsonb
                """
            )

            cursor.execute(
                """
                ALTER TABLE public.document_chunks
                ADD COLUMN IF NOT EXISTS chunk_data
                JSONB DEFAULT '{}'::jsonb
                """
            )

            # ------------------------------------------------
            # INDEXES
            # ------------------------------------------------

            cursor.execute(
                """
                CREATE INDEX IF NOT EXISTS idx_documents_user
                ON public.documents(user_id)
                """
            )

            cursor.execute(
                """
                CREATE INDEX IF NOT EXISTS idx_chunks_document
                ON public.document_chunks(document_id)
                """
            )

            cursor.execute(
                """
                CREATE INDEX IF NOT EXISTS idx_conversations_user
                ON public.conversations(user_id)
                """
            )

            cursor.execute(
                """
                CREATE INDEX IF NOT EXISTS idx_messages_conversation
                ON public.messages(conversation_id)
                """
            )

            cursor.execute(
                """
                CREATE INDEX IF NOT EXISTS idx_chunks_content
                ON public.document_chunks
                USING gin(
                    to_tsvector(
                        'english',
                        COALESCE(content, '')
                    )
                )
                """
            )

            # ------------------------------------------------
            # ADMIN SETUP
            # ------------------------------------------------

            cursor.execute(
                """
                SELECT id
                FROM public.users
                WHERE LOWER(COALESCE(role, 'user')) = 'admin'
                LIMIT 1
                """
            )

            admin = cursor.fetchone()

            if not admin and ADMIN_USERNAME:

                cursor.execute(
                    """
                    UPDATE public.users
                    SET role = 'admin'
                    WHERE LOWER(username) = LOWER(%s)
                    """,
                    (ADMIN_USERNAME,),
                )

            # If no configured admin exists, make the oldest
            # existing user an admin.
            cursor.execute(
                """
                SELECT id
                FROM public.users
                WHERE LOWER(COALESCE(role, 'user')) = 'admin'
                LIMIT 1
                """
            )

            admin = cursor.fetchone()

            if not admin:

                cursor.execute(
                    """
                    SELECT id
                    FROM public.users
                    ORDER BY id
                    LIMIT 1
                    """
                )

                first_user = cursor.fetchone()

                if first_user:

                    cursor.execute(
                        """
                        UPDATE public.users
                        SET role = 'admin'
                        WHERE id = %s
                        """,
                        (first_user[0],),
                    )

            conn.commit()

        except Exception:

            conn.rollback()
            raise

        finally:

            if cursor:
                cursor.close()

    # ========================================================
    # USER LOOKUP
    # ========================================================

    def get_user_by_id(self, user_id):

        conn = self.connect()
        cursor = None

        try:

            cursor = conn.cursor(
                cursor_factory=RealDictCursor
            )

            cursor.execute(
                """
                SELECT
                    id,
                    username,
                    role
                FROM public.users
                WHERE id = %s
                """,
                (int(user_id),),
            )

            row = cursor.fetchone()

            conn.rollback()

            return dict(row) if row else None

        except Exception:

            conn.rollback()
            raise

        finally:

            if cursor:
                cursor.close()

    # ========================================================
    # CREATE USER
    # ========================================================

    def create_user(self, username, password):

        username = str(username or "").strip()

        if not username:
            return None

        if not password:
            return None

        conn = self.connect()
        cursor = None

        try:

            cursor = conn.cursor()

            # First user becomes admin.
            cursor.execute(
                """
                SELECT COUNT(*)
                FROM public.users
                """
            )

            user_count = int(
                cursor.fetchone()[0]
            )

            role = (
                "admin"
                if user_count == 0
                else "user"
            )

            password_hash = generate_password_hash(
                password
            )

            cursor.execute(
                """
                INSERT INTO public.users
                (
                    username,
                    password_hash,
                    role
                )
                VALUES (%s, %s, %s)
                RETURNING id, username, role
                """,
                (
                    username,
                    password_hash,
                    role,
                ),
            )

            result = cursor.fetchone()

            conn.commit()

            return result

        except Exception as exc:

            conn.rollback()

            # Duplicate username.
            if psycopg2 is not None and isinstance(
                exc,
                psycopg2.IntegrityError,
            ):
                return None

            return None

        finally:

            if cursor:
                cursor.close()

    # ========================================================
    # AUTHENTICATE
    # ========================================================

    def authenticate(self, username, password):

        username = str(username or "").strip()

        if not username or not password:
            return None

        conn = self.connect()
        cursor = None

        try:

            cursor = conn.cursor(
                cursor_factory=RealDictCursor
            )

            cursor.execute(
                """
                SELECT
                    id,
                    username,
                    password_hash,
                    role
                FROM public.users
                WHERE LOWER(username) = LOWER(%s)
                """,
                (username,),
            )

            user = cursor.fetchone()

            conn.rollback()

            if not user:
                return None

            stored_hash = user.get(
                "password_hash"
            )

            if not stored_hash:
                return None

            try:

                valid = check_password_hash(
                    stored_hash,
                    password,
                )

            except Exception:

                valid = False

            if not valid:
                return None

            return dict(user)

        except Exception:

            conn.rollback()
            raise

        finally:

            if cursor:
                cursor.close()

    # ========================================================
    # GET USERS
    # ========================================================

    def get_users(self):

        conn = self.connect()
        cursor = None

        try:

            cursor = conn.cursor(
                cursor_factory=RealDictCursor
            )

            cursor.execute(
                """
                SELECT
                    id,
                    username,
                    role,
                    created_at
                FROM public.users
                ORDER BY id
                """
            )

            rows = cursor.fetchall()

            conn.rollback()

            return [
                dict(row)
                for row in rows
            ]

        finally:

            if cursor:
                cursor.close()

    # ========================================================
    # CHANGE USER ROLE
    # ========================================================

    def set_user_role(self, user_id, role):

        role = str(role or "").lower()

        if role not in {"admin", "user"}:
            raise ValueError(
                "Invalid role."
            )

        conn = self.connect()
        cursor = None

        try:

            cursor = conn.cursor()

            cursor.execute(
                """
                SELECT role
                FROM public.users
                WHERE id = %s
                """,
                (int(user_id),),
            )

            current = cursor.fetchone()

            if not current:
                conn.rollback()
                return False

            current_role = str(
                current[0] or "user"
            ).lower()

            if (
                current_role == "admin"
                and role == "user"
            ):

                cursor.execute(
                    """
                    SELECT COUNT(*)
                    FROM public.users
                    WHERE LOWER(
                        COALESCE(role, 'user')
                    ) = 'admin'
                    """
                )

                admin_count = int(
                    cursor.fetchone()[0]
                )

                if admin_count <= 1:
                    conn.rollback()
                    return False

            cursor.execute(
                """
                UPDATE public.users
                SET role = %s
                WHERE id = %s
                """,
                (
                    role,
                    int(user_id),
                ),
            )

            changed = cursor.rowcount > 0

            conn.commit()

            return changed

        except Exception:

            conn.rollback()
            raise

        finally:

            if cursor:
                cursor.close()

    # ========================================================
    # DELETE USER
    # ========================================================

    def delete_user(self, user_id):

        conn = self.connect()
        cursor = None

        try:

            cursor = conn.cursor()

            cursor.execute(
                """
                SELECT role
                FROM public.users
                WHERE id = %s
                """,
                (int(user_id),),
            )

            user = cursor.fetchone()

            if not user:
                conn.rollback()
                return False

            if str(
                user[0] or "user"
            ).lower() == "admin":

                cursor.execute(
                    """
                    SELECT COUNT(*)
                    FROM public.users
                    WHERE LOWER(
                        COALESCE(role, 'user')
                    ) = 'admin'
                    """
                )

                admin_count = int(
                    cursor.fetchone()[0]
                )

                if admin_count <= 1:
                    conn.rollback()
                    return False

            cursor.execute(
                """
                DELETE FROM public.users
                WHERE id = %s
                """,
                (int(user_id),),
            )

            deleted = cursor.rowcount > 0

            conn.commit()

            return deleted

        except Exception:

            conn.rollback()
            raise

        finally:

            if cursor:
                cursor.close()

    # ========================================================
    # RESET PASSWORD
    # ========================================================

    def reset_user_password(
        self,
        user_id,
        password,
    ):

        if not password:
            raise ValueError(
                "Password cannot be empty."
            )

        conn = self.connect()
        cursor = None

        try:

            cursor = conn.cursor()

            password_hash = generate_password_hash(
                password
            )

            cursor.execute(
                """
                UPDATE public.users
                SET password_hash = %s
                WHERE id = %s
                """,
                (
                    password_hash,
                    int(user_id),
                ),
            )

            if cursor.rowcount == 0:
                raise RuntimeError(
                    "User does not exist."
                )

            conn.commit()

        except Exception:

            conn.rollback()
            raise

        finally:

            if cursor:
                cursor.close()

    # ========================================================
    # SAVE DOCUMENT + CHUNKS ATOMICALLY
    # ========================================================

    def save_document_with_chunks(
        self,
        user_id,
        filename,
        file_type,
        file_bytes,
        chunking_method,
        chunk_size,
        chunk_overlap,
        chunks,
        metadata=None,
    ):

        conn = self.connect()
        cursor = None

        try:

            cursor = conn.cursor()

            user_id = int(user_id)

            # ------------------------------------------------
            # Verify owner.
            # ------------------------------------------------

            cursor.execute(
                """
                SELECT id
                FROM public.users
                WHERE id = %s
                FOR SHARE
                """,
                (user_id,),
            )

            if not cursor.fetchone():

                raise RuntimeError(
                    "The logged-in user no longer exists."
                )

            # ------------------------------------------------
            # Remove old version.
            # ------------------------------------------------

            cursor.execute(
                """
                DELETE FROM public.documents
                WHERE user_id = %s
                AND filename = %s
                """,
                (
                    user_id,
                    filename,
                ),
            )

            # ------------------------------------------------
            # Insert document.
            # ------------------------------------------------

            cursor.execute(
                """
                INSERT INTO public.documents
                (
                    user_id,
                    filename,
                    file_type,
                    file_size,
                    file_data,
                    chunking_method,
                    chunk_size,
                    chunk_overlap,
                    metadata
                )
                VALUES
                (
                    %s,
                    %s,
                    %s,
                    %s,
                    %s,
                    %s,
                    %s,
                    %s,
                    %s
                )
                RETURNING id
                """,
                (
                    user_id,
                    filename,
                    file_type,
                    len(file_bytes),
                    psycopg2.Binary(file_bytes),
                    chunking_method,
                    int(chunk_size),
                    int(chunk_overlap),
                    Json(metadata or {}),
                ),
            )

            row = cursor.fetchone()

            if not row:
                raise RuntimeError(
                    "PostgreSQL did not return a document ID."
                )

            document_id = int(row[0])

            # ------------------------------------------------
            # Insert chunks.
            # ------------------------------------------------

            saved_chunks = 0

            for index, original_chunk in enumerate(
                chunks,
                start=1,
            ):

                if not isinstance(
                    original_chunk,
                    dict,
                ):
                    chunk = {
                        "content": str(
                            original_chunk
                        )
                    }
                else:
                    chunk = dict(
                        original_chunk
                    )

                chunk_id = int(
                    chunk.get(
                        "chunk_id",
                        index,
                    )
                )

                content = str(
                    chunk.get(
                        "content",
                        "",
                    )
                    or ""
                )

                if not content.strip():
                    continue

                page = chunk.get("page")

                if page is not None:
                    page = str(page)

                chunk_type = chunk.get(
                    "chunk_type",
                    "Document",
                )

                chunk_metadata = chunk.get(
                    "metadata",
                    {},
                )

                if not isinstance(
                    chunk_metadata,
                    dict,
                ):
                    chunk_metadata = {}

                tokens = chunk.get("tokens")

                if tokens is None:
                    tokens = max(
                        1,
                        len(
                            content.split()
                        ),
                    )

                characters = chunk.get(
                    "characters"
                )

                if characters is None:
                    characters = len(content)

                cursor.execute(
                    """
                    INSERT INTO public.document_chunks
                    (
                        document_id,
                        chunk_id,
                        chunk_type,
                        content,
                        page,
                        tokens,
                        characters,
                        metadata,
                        chunk_data
                    )
                    VALUES
                    (
                        %s,
                        %s,
                        %s,
                        %s,
                        %s,
                        %s,
                        %s,
                        %s,
                        %s
                    )
                    ON CONFLICT
                    (
                        document_id,
                        chunk_id
                    )
                    DO UPDATE SET
                        chunk_type = EXCLUDED.chunk_type,
                        content = EXCLUDED.content,
                        page = EXCLUDED.page,
                        tokens = EXCLUDED.tokens,
                        characters = EXCLUDED.characters,
                        metadata = EXCLUDED.metadata,
                        chunk_data = EXCLUDED.chunk_data
                    """,
                    (
                        document_id,
                        chunk_id,
                        chunk_type,
                        content,
                        page,
                        int(tokens),
                        int(characters),
                        Json(chunk_metadata),
                        Json(chunk),
                    ),
                )

                saved_chunks += 1

            if saved_chunks == 0:

                raise RuntimeError(
                    "No usable document chunks were generated."
                )

            # ------------------------------------------------
            # Verify chunks.
            # ------------------------------------------------

            cursor.execute(
                """
                SELECT COUNT(*)
                FROM public.document_chunks
                WHERE document_id = %s
                """,
                (document_id,),
            )

            chunk_count = int(
                cursor.fetchone()[0]
            )

            if chunk_count != saved_chunks:

                raise RuntimeError(
                    "Chunk verification failed. "
                    f"Expected {saved_chunks}, "
                    f"found {chunk_count}."
                )

            # ------------------------------------------------
            # Commit.
            # ------------------------------------------------

            conn.commit()

            return {
                "document_id": document_id,
                "chunk_count": chunk_count,
            }

        except Exception:

            conn.rollback()
            raise

        finally:

            if cursor:
                cursor.close()

    # ========================================================
    # GET DOCUMENTS
    # ========================================================

    def get_documents(
        self,
        user_id=None,
        is_admin=False,
        search=None,
    ):

        conn = self.connect()
        cursor = None

        try:

            cursor = conn.cursor(
                cursor_factory=RealDictCursor
            )

            conditions = []
            params = []

            if not is_admin:

                if user_id is None:
                    return []

                conditions.append(
                    "d.user_id = %s"
                )

                params.append(
                    int(user_id)
                )

            if search and str(search).strip():

                pattern = (
                    "%"
                    + str(search).strip()
                    + "%"
                )

                conditions.append(
                    """
                    (
                        LOWER(d.filename)
                        LIKE LOWER(%s)
                        OR LOWER(
                            COALESCE(
                                d.file_type,
                                ''
                            )
                        )
                        LIKE LOWER(%s)
                    )
                    """
                )

                params.extend(
                    [
                        pattern,
                        pattern,
                    ]
                )

            where_clause = ""

            if conditions:
                where_clause = (
                    "WHERE "
                    + " AND ".join(conditions)
                )

            query = f"""
                SELECT
                    d.id,
                    d.user_id,
                    d.filename,
                    d.file_type,
                    d.file_size,
                    d.chunking_method,
                    d.chunk_size,
                    d.chunk_overlap,
                    d.metadata,
                    d.created_at,
                    u.username,

                    (
                        SELECT COUNT(*)
                        FROM public.document_chunks c
                        WHERE c.document_id = d.id
                    ) AS chunk_count

                FROM public.documents d

                LEFT JOIN public.users u
                    ON u.id = d.user_id

                {where_clause}

                ORDER BY
                    d.created_at DESC,
                    d.id DESC
            """

            cursor.execute(
                query,
                params,
            )

            rows = cursor.fetchall()

            conn.rollback()

            return [
                dict(row)
                for row in rows
            ]

        except Exception:

            conn.rollback()
            raise

        finally:

            if cursor:
                cursor.close()

    # ========================================================
    # DOCUMENT BY ID
    # ========================================================

    def get_document_by_id(
        self,
        document_id,
        user_id=None,
        is_admin=False,
    ):

        conn = self.connect()
        cursor = None

        try:

            cursor = conn.cursor(
                cursor_factory=RealDictCursor
            )

            if is_admin:

                cursor.execute(
                    """
                    SELECT
                        d.*,
                        u.username
                    FROM public.documents d
                    LEFT JOIN public.users u
                        ON u.id = d.user_id
                    WHERE d.id = %s
                    """,
                    (int(document_id),),
                )

            else:

                cursor.execute(
                    """
                    SELECT
                        d.*,
                        u.username
                    FROM public.documents d
                    LEFT JOIN public.users u
                        ON u.id = d.user_id
                    WHERE d.id = %s
                    AND d.user_id = %s
                    """,
                    (
                        int(document_id),
                        int(user_id),
                    ),
                )

            row = cursor.fetchone()

            conn.rollback()

            return dict(row) if row else None

        except Exception:

            conn.rollback()
            raise

        finally:

            if cursor:
                cursor.close()

    # ========================================================
    # DELETE DOCUMENT
    # ========================================================

    def delete_document(
        self,
        document_id,
        user_id=None,
        is_admin=False,
    ):

        conn = self.connect()
        cursor = None

        try:

            cursor = conn.cursor()

            if is_admin:

                cursor.execute(
                    """
                    DELETE FROM public.documents
                    WHERE id = %s
                    """,
                    (int(document_id),),
                )

            else:

                cursor.execute(
                    """
                    DELETE FROM public.documents
                    WHERE id = %s
                    AND user_id = %s
                    """,
                    (
                        int(document_id),
                        int(user_id),
                    ),
                )

            deleted = cursor.rowcount > 0

            conn.commit()

            return deleted

        except Exception:

            conn.rollback()
            raise

        finally:

            if cursor:
                cursor.close()

    # ========================================================
    # SEARCH CHUNKS
    # ========================================================

    def search_chunks(
        self,
        user_id,
        query,
        selected_document_ids=None,
        chunk_types=None,
        top_k=5,
    ):

        query = str(query or "").strip()

        if not query:
            return []

        conn = self.connect()
        cursor = None

        try:

            cursor = conn.cursor(
                cursor_factory=RealDictCursor
            )

            # ------------------------------------------------
            # SECURITY CONDITIONS
            # ------------------------------------------------

            conditions = [
                "d.user_id = %s"
            ]

            condition_params = [
                int(user_id)
            ]

            # ------------------------------------------------
            # DOCUMENT FILTER
            # ------------------------------------------------

            if selected_document_ids:

                clean_ids = []

                for value in selected_document_ids:

                    try:
                        clean_ids.append(
                            int(value)
                        )
                    except Exception:
                        continue

                if clean_ids:

                    conditions.append(
                        "d.id = ANY(%s)"
                    )

                    condition_params.append(
                        clean_ids
                    )

            # ------------------------------------------------
            # CHUNK TYPE FILTER
            # ------------------------------------------------

            if chunk_types:

                clean_types = [
                    str(x)
                    for x in chunk_types
                    if str(x).strip()
                ]

                if clean_types:

                    conditions.append(
                        "c.chunk_type = ANY(%s)"
                    )

                    condition_params.append(
                        clean_types
                    )

            where_clause = " AND ".join(
                conditions
            )

            # ------------------------------------------------
            # IMPORTANT FIX:
            #
            # The placeholders in SQL occur in this order:
            #
            # 1. similarity query
            # 2. user_id / filters
            # 3. tsquery
            # 4. LIKE content
            # 5. LIKE filename
            # 6. LIMIT
            #
            # The original app.py put user_id first.
            # ------------------------------------------------

            sql = f"""
                SELECT
                    c.id,
                    c.document_id,
                    c.chunk_id,
                    c.chunk_type,
                    c.content,
                    c.page,
                    c.tokens,
                    c.characters,
                    c.metadata,
                    c.chunk_data,
                    c.created_at,

                    d.filename,
                    d.file_type,
                    d.chunking_method,

                    ts_rank(
                        to_tsvector(
                            'english',
                            COALESCE(c.content, '')
                        ),
                        plainto_tsquery(
                            'english',
                            %s
                        )
                    ) AS similarity_score

                FROM public.document_chunks c

                INNER JOIN public.documents d
                    ON d.id = c.document_id

                WHERE
                    {where_clause}

                    AND
                    (
                        to_tsvector(
                            'english',
                            COALESCE(c.content, '')
                        )
                        @@ plainto_tsquery(
                            'english',
                            %s
                        )

                        OR LOWER(
                            COALESCE(c.content, '')
                        )
                        LIKE LOWER(%s)

                        OR LOWER(
                            COALESCE(d.filename, '')
                        )
                        LIKE LOWER(%s)
                    )

                ORDER BY
                    similarity_score DESC,
                    c.id DESC

                LIMIT %s
            """

            params = [
                query,
                *condition_params,
                query,
                f"%{query}%",
                f"%{query}%",
                max(
                    1,
                    int(top_k),
                ),
            ]

            cursor.execute(
                sql,
                params,
            )

            rows = cursor.fetchall()

            conn.rollback()

            return [
                dict(row)
                for row in rows
            ]

        except Exception:

            conn.rollback()
            raise

        finally:

            if cursor:
                cursor.close()

    # ========================================================
    # CONVERSATIONS
    # ========================================================

    def create_conversation(
        self,
        user_id,
        title="New Conversation",
    ):

        conn = self.connect()
        cursor = None

        try:

            cursor = conn.cursor()

            cursor.execute(
                """
                SELECT id
                FROM public.users
                WHERE id = %s
                """,
                (int(user_id),),
            )

            if not cursor.fetchone():

                raise RuntimeError(
                    "User no longer exists."
                )

            cursor.execute(
                """
                INSERT INTO public.conversations
                (
                    user_id,
                    title
                )
                VALUES
                (
                    %s,
                    %s
                )
                RETURNING id
                """,
                (
                    int(user_id),
                    str(title),
                ),
            )

            conversation_id = int(
                cursor.fetchone()[0]
            )

            conn.commit()

            return conversation_id

        except Exception:

            conn.rollback()
            raise

        finally:

            if cursor:
                cursor.close()

    # ========================================================
    # GET CONVERSATIONS
    # ========================================================

    def get_conversations(self, user_id):

        conn = self.connect()
        cursor = None

        try:

            cursor = conn.cursor(
                cursor_factory=RealDictCursor
            )

            cursor.execute(
                """
                SELECT
                    id,
                    user_id,
                    title,
                    created_at
                FROM public.conversations
                WHERE user_id = %s
                ORDER BY
                    created_at DESC,
                    id DESC
                """,
                (int(user_id),),
            )

            rows = cursor.fetchall()

            conn.rollback()

            return [
                dict(row)
                for row in rows
            ]

        except Exception:

            conn.rollback()
            raise

        finally:

            if cursor:
                cursor.close()

    # ========================================================
    # GET MESSAGES
    # ========================================================

    def get_messages(
        self,
        conversation_id,
        user_id=None,
    ):

        conn = self.connect()
        cursor = None

        try:

            cursor = conn.cursor(
                cursor_factory=RealDictCursor
            )

            if user_id is None:

                cursor.execute(
                    """
                    SELECT
                        m.id,
                        m.conversation_id,
                        m.role,
                        m.content,
                        m.created_at
                    FROM public.messages m
                    WHERE m.conversation_id = %s
                    ORDER BY m.id ASC
                    """,
                    (int(conversation_id),),
                )

            else:

                cursor.execute(
                    """
                    SELECT
                        m.id,
                        m.conversation_id,
                        m.role,
                        m.content,
                        m.created_at
                    FROM public.messages m
                    INNER JOIN public.conversations c
                        ON c.id = m.conversation_id
                    WHERE m.conversation_id = %s
                    AND c.user_id = %s
                    ORDER BY m.id ASC
                    """,
                    (
                        int(conversation_id),
                        int(user_id),
                    ),
                )

            rows = cursor.fetchall()

            conn.rollback()

            return [
                dict(row)
                for row in rows
            ]

        except Exception:

            conn.rollback()
            raise

        finally:

            if cursor:
                cursor.close()

    # ========================================================
    # SAVE MESSAGE
    # ========================================================

    def save_message(
        self,
        conversation_id,
        role,
        content,
        user_id=None,
    ):

        if not content:
            return

        conn = self.connect()
        cursor = None

        try:

            cursor = conn.cursor()

            if user_id is not None:

                cursor.execute(
                    """
                    SELECT id
                    FROM public.conversations
                    WHERE id = %s
                    AND user_id = %s
                    """,
                    (
                        int(conversation_id),
                        int(user_id),
                    ),
                )

                if not cursor.fetchone():

                    raise RuntimeError(
                        "Conversation does not belong "
                        "to the current user."
                    )

            cursor.execute(
                """
                INSERT INTO public.messages
                (
                    conversation_id,
                    role,
                    content
                )
                VALUES
                (
                    %s,
                    %s,
                    %s
                )
                """,
                (
                    int(conversation_id),
                    str(role),
                    str(content),
                ),
            )

            conn.commit()

        except Exception:

            conn.rollback()
            raise

        finally:

            if cursor:
                cursor.close()

    # ========================================================
    # UPDATE CONVERSATION TITLE
    # ========================================================

    def update_conversation_title(
        self,
        conversation_id,
        title,
        user_id=None,
    ):

        conn = self.connect()
        cursor = None

        try:

            cursor = conn.cursor()

            if user_id is None:

                cursor.execute(
                    """
                    UPDATE public.conversations
                    SET title = %s
                    WHERE id = %s
                    """,
                    (
                        str(title),
                        int(conversation_id),
                    ),
                )

            else:

                cursor.execute(
                    """
                    UPDATE public.conversations
                    SET title = %s
                    WHERE id = %s
                    AND user_id = %s
                    """,
                    (
                        str(title),
                        int(conversation_id),
                        int(user_id),
                    ),
                )

            conn.commit()

        except Exception:

            conn.rollback()
            raise

        finally:

            if cursor:
                cursor.close()

    # ========================================================
    # STATISTICS
    # ========================================================

    def get_stats(self):

        conn = self.connect()
        cursor = None

        try:

            cursor = conn.cursor()

            stats = {}

            for key, table in [
                ("users", "users"),
                ("documents", "documents"),
                ("chunks", "document_chunks"),
                ("conversations", "conversations"),
                ("messages", "messages"),
            ]:

                cursor.execute(
                    f"""
                    SELECT COUNT(*)
                    FROM public.{table}
                    """
                )

                stats[key] = int(
                    cursor.fetchone()[0]
                )

            conn.rollback()

            return stats

        except Exception:

            conn.rollback()
            raise

        finally:

            if cursor:
                cursor.close()


# ============================================================
# DOCUMENT EXTRACTION
# ============================================================

def extract_document(
    filename,
    file_bytes,
):

    extension = Path(
        filename
    ).suffix.lower()

    if extension not in SUPPORTED_EXTENSIONS:

        raise ValueError(
            f"Unsupported file type: {extension}"
        )

    # ========================================================
    # PDF
    # ========================================================

    if extension == ".pdf":

        if PdfReader is None:

            raise RuntimeError(
                "pypdf is not installed."
            )

        reader = PdfReader(
            io.BytesIO(file_bytes)
        )

        pages = []

        for page_number, page in enumerate(
            reader.pages,
            start=1,
        ):

            text = page.extract_text() or ""

            if text.strip():

                pages.append(
                    {
                        "page": page_number,
                        "content": text,
                        "chunk_type": "Document",
                        "metadata": {
                            "page": page_number
                        },
                    }
                )

        return pages

    # ========================================================
    # DOCX
    # ========================================================

    if extension == ".docx":

        if DocxDocument is None:

            raise RuntimeError(
                "python-docx is not installed."
            )

        document = DocxDocument(
            io.BytesIO(file_bytes)
        )

        parts = []

        for paragraph in document.paragraphs:

            text = str(
                paragraph.text or ""
            ).strip()

            if text:
                parts.append(text)

        # Also include table content.
        for table_index, table in enumerate(
            document.tables,
            start=1,
        ):

            rows = []

            for row in table.rows:

                rows.append(
                    " | ".join(
                        str(
                            cell.text or ""
                        ).strip()
                        for cell in row.cells
                    )
                )

            if rows:

                parts.append(
                    "\n".join(
                        [
                            f"TABLE {table_index}",
                            *rows,
                        ]
                    )
                )

        content = "\n\n".join(parts)

        if not content.strip():
            return []

        return [
            {
                "page": None,
                "content": content,
                "chunk_type": "Document",
            }
        ]

    # ========================================================
    # PPTX
    # ========================================================

    if extension == ".pptx":

        if Presentation is None:

            raise RuntimeError(
                "python-pptx is not installed."
            )

        presentation = Presentation(
            io.BytesIO(file_bytes)
        )

        pages = []

        for slide_number, slide in enumerate(
            presentation.slides,
            start=1,
        ):

            parts = []

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    text = str(
                        shape.text or ""
                    ).strip()

                    if text:
                        parts.append(text)

            content = "\n".join(parts)

            if content.strip():

                pages.append(
                    {
                        "page": slide_number,
                        "content": content,
                        "chunk_type": "Document",
                        "metadata": {
                            "slide": slide_number
                        },
                    }
                )

        return pages

    # ========================================================
    # XLSX / XLS
    # ========================================================

    if extension in {".xlsx", ".xls"}:

        try:

            sheets = pd.read_excel(
                io.BytesIO(file_bytes),
                sheet_name=None,
            )

        except ImportError as exc:

            raise RuntimeError(
                "Excel support requires the appropriate "
                "pandas Excel engine. For .xlsx install "
                "openpyxl; for .xls install xlrd."
            ) from exc

        except Exception as exc:

            raise RuntimeError(
                f"Could not read Excel file: {exc}"
            ) from exc

        result = []

        for sheet_name, dataframe in sheets.items():

            dataframe = dataframe.fillna("")

            content = dataframe.to_csv(
                index=False
            )

            if not content.strip():
                continue

            result.append(
                {
                    "page": str(sheet_name),
                    "content": content,
                    "chunk_type": "Table",
                    "metadata": {
                        "sheet": str(sheet_name)
                    },
                }
            )

        return result

    # ========================================================
    # CSV
    # ========================================================

    if extension == ".csv":

        try:

            dataframe = pd.read_csv(
                io.BytesIO(file_bytes)
            )

        except Exception:

            # Fallback for non-standard CSV encoding.
            dataframe = pd.read_csv(
                io.BytesIO(file_bytes),
                encoding="latin-1",
            )

        dataframe = dataframe.fillna("")

        content = dataframe.to_csv(
            index=False
        )

        return [
            {
                "page": None,
                "content": content,
                "chunk_type": "Table",
            }
        ]

    # ========================================================
    # TXT / MARKDOWN
    # ========================================================

    content = file_bytes.decode(
        "utf-8",
        errors="replace",
    )

    if not content.strip():
        return []

    return [
        {
            "page": None,
            "content": content,
            "chunk_type": (
                "Markdown"
                if extension in {
                    ".md",
                    ".markdown",
                }
                else "Document"
            ),
        }
    ]


# ============================================================
# CHUNK DOCUMENT
# ============================================================

def chunk_document(
    extracted_pages,
    method,
    chunk_size,
    chunk_overlap,
    filename,
):

    if not extracted_pages:
        return []

    all_chunks = []

    for page_data in extracted_pages:

        if not isinstance(
            page_data,
            dict,
        ):
            continue

        content = str(
            page_data.get(
                "content",
                "",
            )
            or ""
        )

        if not content.strip():
            continue

        page = page_data.get("page")

        chunk_type = page_data.get(
            "chunk_type",
            "Document",
        )

        metadata = page_data.get(
            "metadata",
            {},
        )

        if not isinstance(
            metadata,
            dict,
        ):
            metadata = {}

        tables = []

        # ----------------------------------------------------
        # Tables
        # ----------------------------------------------------

        if chunk_type == "Table":

            try:

                dataframe = pd.read_csv(
                    io.StringIO(content)
                )

                dataframe = dataframe.fillna("")

                tables.append(
                    {
                        "table": dataframe.to_dict(
                            orient="records"
                        ),
                        "page": page,
                        "sheet": metadata.get(
                            "sheet"
                        ),
                        "source": filename,
                    }
                )

            except Exception:
                pass

        engine = ChunkingEngine(
            text=content,
            metadata={
                **metadata,
                "filename": filename,
            },
            tables=tables,
            source=filename,
            chunk_size=int(chunk_size),
            chunk_overlap=int(chunk_overlap),
            pages=[
                {
                    "page": page,
                    "text": content,
                }
            ],
        )

        try:

            # ------------------------------------------------
            # Tables must use table chunking.
            # ------------------------------------------------

            if chunk_type == "Table":

                chunks = engine.chunk(
                    "Table"
                )

            elif method == "Multimodal":

                chunks = engine.chunk(
                    "Multimodal"
                )

            else:

                chunks = engine.chunk(
                    method
                )

        except Exception as exc:

            st.warning(
                f"Chunking method '{method}' failed "
                f"for {filename}: {exc}. "
                "Falling back to Recursive."
            )

            chunks = engine.chunk(
                "Recursive"
            )

        for chunk in chunks:

            if not isinstance(
                chunk,
                dict,
            ):
                continue

            chunk = dict(chunk)

            chunk["filename"] = filename

            if chunk.get("page") is None:
                chunk["page"] = page

            chunk.setdefault(
                "chunk_type",
                chunk_type,
            )

            content_value = str(
                chunk.get(
                    "content",
                    "",
                )
                or ""
            )

            if not content_value.strip():
                continue

            chunk["content"] = content_value

            chunk.setdefault(
                "characters",
                len(content_value),
            )

            chunk.setdefault(
                "tokens",
                max(
                    1,
                    len(
                        content_value.split()
                    ),
                ),
            )

            chunk.setdefault(
                "metadata",
                {},
            )

            all_chunks.append(chunk)

    # --------------------------------------------------------
    # Global chunk IDs.
    # --------------------------------------------------------

    for index, chunk in enumerate(
        all_chunks,
        start=1,
    ):

        chunk["chunk_id"] = index

    return all_chunks


# ============================================================
# LLM
# ============================================================

@st.cache_resource
def create_llm():

    if not OPENROUTER_API_KEY:
        return None

    try:

        llm = GeminiModel()

        return llm

    except Exception:

        return None


# ============================================================
# SESSION STATE
# ============================================================

def initialize_session_state():

    defaults = {
        "authenticated": False,
        "user": None,
        "db": None,
        "llm": None,
        "conversation_id": None,
        "selected_document_ids": [],
        "selected_chunk_types": [],
        "top_k": DEFAULT_TOP_K,
        "chunking_method": "Recursive",
        "document_search": "",
    }

    for key, value in defaults.items():

        if key not in st.session_state:

            st.session_state[key] = value


initialize_session_state()


# ============================================================
# DATABASE STARTUP
# ============================================================

if st.session_state.db is None:

    try:

        db = PostgreSQLStore()

        db.initialize()

        st.session_state.db = db

    except Exception as exc:

        st.error(
            "❌ Database initialization failed."
        )

        st.exception(exc)

        st.markdown(
            """
### PostgreSQL configuration

Configure either:

`DATABASE_URL`

or:

- `POSTGRES_HOST`
- `POSTGRES_PORT`
- `POSTGRES_DB`
- `POSTGRES_USER`
- `POSTGRES_PASSWORD`
"""
        )

        st.stop()


# ============================================================
# USER VALIDATION
# ============================================================

def validate_logged_in_user():

    if not st.session_state.authenticated:
        return False

    user = st.session_state.user

    if not user:
        return False

    try:

        fresh_user = (
            st.session_state.db.get_user_by_id(
                user["id"]
            )
        )

        if not fresh_user:

            st.session_state.authenticated = False
            st.session_state.user = None
            st.session_state.conversation_id = None
            st.session_state.selected_document_ids = []

            return False

        st.session_state.user = fresh_user

        return True

    except Exception:

        return False


# ============================================================
# LOGIN PAGE
# ============================================================

def login_page():

    st.title(
        "📄 Enterprise Document Intelligence"
    )

    st.caption(
        "PostgreSQL • RAG • OpenRouter • Gemini"
    )

    login_tab, register_tab = st.tabs(
        [
            "🔐 Login",
            "📝 Register",
        ]
    )

    # ========================================================
    # LOGIN
    # ========================================================

    with login_tab:

        with st.form(
            "login_form",
            clear_on_submit=False,
        ):

            username = st.text_input(
                "Username"
            )

            password = st.text_input(
                "Password",
                type="password",
            )

            submitted = st.form_submit_button(
                "Login",
                use_container_width=True,
            )

            if submitted:

                if not username or not password:

                    st.error(
                        "Enter username and password."
                    )

                else:

                    try:

                        user = (
                            st.session_state.db.authenticate(
                                username,
                                password,
                            )
                        )

                        if user:

                            st.session_state.authenticated = True
                            st.session_state.user = user
                            st.session_state.conversation_id = None
                            st.session_state.selected_document_ids = []
                            st.session_state.selected_chunk_types = []
                            st.session_state.llm = create_llm()

                            st.rerun()

                        else:

                            st.error(
                                "Invalid username or password."
                            )

                    except Exception as exc:

                        st.error(
                            "Login failed."
                        )

                        st.exception(exc)

    # ========================================================
    # REGISTER
    # ========================================================

    with register_tab:

        st.caption(
            "The first registered account becomes an administrator."
        )

        with st.form(
            "register_form",
            clear_on_submit=False,
        ):

            username = st.text_input(
                "New username"
            )

            password = st.text_input(
                "Password",
                type="password",
            )

            confirm = st.text_input(
                "Confirm password",
                type="password",
            )

            submitted = st.form_submit_button(
                "Create account",
                use_container_width=True,
            )

            if submitted:

                if not username.strip():

                    st.error(
                        "Username is required."
                    )

                elif len(password) < 6:

                    st.error(
                        "Password must contain at least 6 characters."
                    )

                elif password != confirm:

                    st.error(
                        "Passwords do not match."
                    )

                else:

                    result = (
                        st.session_state.db.create_user(
                            username,
                            password,
                        )
                    )

                    if result:

                        role = result[2]

                        if role == "admin":

                            st.success(
                                "Account created. "
                                "This is the first account, "
                                "so it has administrator access."
                            )

                        else:

                            st.success(
                                "Account created. "
                                "You can now log in."
                            )

                    else:

                        st.error(
                            "Could not create account. "
                            "The username may already exist."
                        )


# ============================================================
# LOGOUT
# ============================================================

def logout():

    st.session_state.authenticated = False
    st.session_state.user = None
    st.session_state.conversation_id = None
    st.session_state.llm = None
    st.session_state.selected_document_ids = []
    st.session_state.selected_chunk_types = []

    st.rerun()


# ============================================================
# CHAT HELPERS
# ============================================================

def create_new_conversation():

    user = st.session_state.user

    conversation_id = (
        st.session_state.db.create_conversation(
            user["id"],
            "New Conversation",
        )
    )

    st.session_state.conversation_id = (
        conversation_id
    )


def ensure_conversation():

    if st.session_state.conversation_id:
        return

    create_new_conversation()


def load_conversation_messages():

    conversation_id = (
        st.session_state.conversation_id
    )

    if not conversation_id:
        return []

    return (
        st.session_state.db.get_messages(
            conversation_id,
            user_id=st.session_state.user["id"],
        )
    )


# ============================================================
# RETRIEVAL
# ============================================================

def retrieve_documents(
    query,
    chunk_types=None,
):

    user = st.session_state.user

    return (
        st.session_state.db.search_chunks(
            user_id=user["id"],
            query=query,
            selected_document_ids=(
                st.session_state.selected_document_ids
                or None
            ),
            chunk_types=(
                chunk_types
                if chunk_types is not None
                else (
                    st.session_state.selected_chunk_types
                    or None
                )
            ),
            top_k=st.session_state.top_k,
        )
    )


# ============================================================
# RAG
# ============================================================

def perform_rag(query):

    llm = st.session_state.llm

    # --------------------------------------------------------
    # No LLM.
    # --------------------------------------------------------

    if llm is None:

        local_results = retrieve_documents(query)

        if local_results:

            first = local_results[0]

            return (
                "OpenRouter is not configured, but I found "
                "matching document evidence:\n\n"
                + str(
                    first.get(
                        "content",
                        "",
                    )
                ),
                local_results,
                "documents",
            )

        return (
            "No usable evidence was found. "
            "Configure OPENROUTER_API_KEY to enable "
            "AI-generated answers.",
            [],
            "error",
        )

    # --------------------------------------------------------
    # Ask model which action to use.
    # --------------------------------------------------------

    try:

        plan = llm.plan_action(
            query=query,
            previous_actions=[],
            previous_evaluations=[],
        )

    except Exception:

        plan = {
            "action": "vector_search",
            "query": query,
            "reason": "Default document retrieval.",
        }

    action = plan.get(
        "action",
        "vector_search",
    )

    search_query = str(
        plan.get(
            "query",
            query,
        )
        or query
    ).strip()

    if not search_query:
        search_query = query

    # --------------------------------------------------------
    # VECTOR SEARCH
    # --------------------------------------------------------

    if action == "vector_search":

        local_results = retrieve_documents(
            search_query
        )

        if local_results:

            try:

                evaluation = (
                    llm.evaluate_evidence(
                        query=query,
                        results=local_results,
                        action="vector_search",
                    )
                )

            except Exception:

                evaluation = {
                    "sufficient": True,
                    "confidence": 0.5,
                }

            if evaluation.get(
                "sufficient",
                False,
            ):

                answer = llm.generate_answer(
                    query=query,
                    chunks=local_results,
                    source_type="uploaded documents",
                )

                return (
                    answer,
                    local_results,
                    "documents",
                )

        # Continue to web fallback.
        local_fallback = local_results

    # --------------------------------------------------------
    # DOCUMENT SEARCH
    # --------------------------------------------------------

    elif action == "document_search":

        local_results = retrieve_documents(
            search_query
        )

        if local_results:

            answer = llm.generate_answer(
                query=query,
                chunks=local_results,
                source_type="uploaded documents",
            )

            return (
                answer,
                local_results,
                "documents",
            )

        local_fallback = []

    # --------------------------------------------------------
    # TABLE SEARCH
    # --------------------------------------------------------

    elif action == "table_search":

        local_results = retrieve_documents(
            search_query,
            chunk_types=["Table"],
        )

        if local_results:

            answer = llm.generate_answer(
                query=query,
                chunks=local_results,
                source_type="uploaded tables",
            )

            return (
                answer,
                local_results,
                "tables",
            )

        local_fallback = []

    # --------------------------------------------------------
    # WEB SEARCH
    # --------------------------------------------------------

    elif action == "web_search":

        try:

            web_results = llm.web_search(
                search_query,
                max_results=5,
            )

        except Exception:

            web_results = []

        if web_results:

            answer = llm.generate_answer(
                query=query,
                chunks=web_results,
                source_type="web",
            )

            return (
                answer,
                web_results,
                "web",
            )

        local_fallback = []

    else:

        local_fallback = []

    # --------------------------------------------------------
    # If the model selected a document action but evidence
    # wasn't enough, try local vector search once.
    # --------------------------------------------------------

    if action != "vector_search":

        try:

            local_results = retrieve_documents(
                query
            )

        except Exception:

            local_results = []

    else:

        local_results = local_fallback

    # --------------------------------------------------------
    # Evaluate fallback local evidence.
    # --------------------------------------------------------

    if local_results:

        try:

            evaluation = (
                llm.evaluate_evidence(
                    query=query,
                    results=local_results,
                    action="vector_search",
                )
            )

        except Exception:

            evaluation = {
                "sufficient": False,
                "confidence": 0.0,
            }

        if evaluation.get(
            "sufficient",
            False,
        ):

            answer = llm.generate_answer(
                query=query,
                chunks=local_results,
                source_type="uploaded documents",
            )

            return (
                answer,
                local_results,
                "documents",
            )

    # --------------------------------------------------------
    # Web fallback.
    # --------------------------------------------------------

    try:

        web_results = llm.web_search(
            query,
            max_results=5,
        )

    except Exception:

        web_results = []

    if web_results:

        answer = llm.generate_answer(
            query=query,
            chunks=web_results,
            source_type="web",
        )

        return (
            answer,
            web_results,
            "web",
        )

    # --------------------------------------------------------
    # Weak local evidence.
    # --------------------------------------------------------

    if local_results:

        answer = llm.generate_answer(
            query=query,
            chunks=local_results,
            source_type="uploaded documents",
        )

        return (
            answer,
            local_results,
            "documents",
        )

    # --------------------------------------------------------
    # Direct LLM.
    # --------------------------------------------------------

    try:

        answer = llm.generate(
            query
        )

        return (
            answer,
            [],
            "llm",
        )

    except Exception as exc:

        return (
            f"LLM request failed: {exc}",
            [],
            "error",
        )


# ============================================================
# SIDEBAR
# ============================================================

def render_sidebar():

    user = st.session_state.user

    with st.sidebar:

        st.markdown(
            "## 📄 Enterprise RAG"
        )

        st.caption(
            f"Logged in as **{user['username']}**"
        )

        st.caption(
            f"Role: **{user['role']}**"
        )

        st.divider()

        pages = [
            "💬 Chat",
            "📚 Documents",
        ]

        if str(
            user["role"]
        ).lower() == "admin":

            pages.append(
                "🛠️ Admin"
            )

        page = st.radio(
            "Navigation",
            pages,
        )

        st.divider()

        st.markdown(
            "### Retrieval settings"
        )

        st.session_state.top_k = st.slider(
            "Top results",
            min_value=1,
            max_value=15,
            value=int(
                st.session_state.top_k
            ),
        )

        methods = [
            "Character",
            "Recursive",
            "Token",
            "Markdown",
            "Context",
            "Multimodal",
        ]

        current_method = (
            st.session_state.chunking_method
        )

        if current_method not in methods:
            current_method = "Recursive"

        st.session_state.chunking_method = (
            st.selectbox(
                "Default chunking method",
                methods,
                index=methods.index(
                    current_method
                ),
            )
        )

        st.divider()

        if st.button(
            "➕ New conversation",
            use_container_width=True,
        ):

            create_new_conversation()
            st.rerun()

        if st.button(
            "🚪 Logout",
            use_container_width=True,
        ):

            logout()

        return page


# ============================================================
# CHAT PAGE
# ============================================================

def chat_page():

    st.title(
        "💬 Enterprise Document Assistant"
    )

    if st.session_state.llm is None:

        st.warning(
            "⚠️ OpenRouter is not configured. "
            "Set OPENROUTER_API_KEY in your "
            "environment/configuration."
        )

    selected_ids = (
        st.session_state.selected_document_ids
    )

    if selected_ids:

        docs = []

        for document_id in selected_ids:

            try:

                document = (
                    st.session_state.db.get_document_by_id(
                        document_id,
                        user_id=st.session_state.user["id"],
                        is_admin=False,
                    )
                )

            except Exception:

                document = None

            if document:
                docs.append(document)

        if docs:

            st.info(
                "🎯 Searching only selected document(s): "
                + ", ".join(
                    d["filename"]
                    for d in docs
                )
            )

        else:

            st.session_state.selected_document_ids = []

    else:

        st.caption(
            "🔎 Chat searches across all your documents."
        )

    ensure_conversation()

    # --------------------------------------------------------
    # Conversations
    # --------------------------------------------------------

    try:

        conversations = (
            st.session_state.db.get_conversations(
                st.session_state.user["id"]
            )
        )

    except Exception as exc:

        st.error(
            "Could not load conversations."
        )

        st.exception(exc)

        return

    if conversations:

        ids = [
            item["id"]
            for item in conversations
        ]

        titles = [
            f"{item['title']} (#{item['id']})"
            for item in conversations
        ]

        current_index = 0

        if (
            st.session_state.conversation_id
            in ids
        ):

            current_index = ids.index(
                st.session_state.conversation_id
            )

        selected = st.selectbox(
            "Conversation",
            titles,
            index=current_index,
        )

        selected_index = titles.index(
            selected
        )

        selected_id = ids[
            selected_index
        ]

        if (
            selected_id
            != st.session_state.conversation_id
        ):

            st.session_state.conversation_id = (
                selected_id
            )

            st.rerun()

    # --------------------------------------------------------
    # Messages
    # --------------------------------------------------------

    messages = load_conversation_messages()

    for message in messages:

        role = message.get("role")

        if role not in {
            "user",
            "assistant",
        }:
            continue

        with st.chat_message(role):

            st.markdown(
                message.get(
                    "content",
                    "",
                )
            )

    # --------------------------------------------------------
    # Prompt
    # --------------------------------------------------------

    prompt = st.chat_input(
        "Ask something about your documents..."
    )

    if not prompt:
        return

    prompt = prompt.strip()

    if not prompt:
        return

    conversation_id = (
        st.session_state.conversation_id
    )

    user_id = st.session_state.user["id"]

    # --------------------------------------------------------
    # Save user message.
    # --------------------------------------------------------

    try:

        st.session_state.db.save_message(
            conversation_id,
            "user",
            prompt,
            user_id=user_id,
        )

    except Exception as exc:

        st.error(
            f"Could not save your message: {exc}"
        )

        return

    with st.chat_message("user"):

        st.markdown(prompt)

    # --------------------------------------------------------
    # Generate answer.
    # --------------------------------------------------------

    with st.chat_message("assistant"):

        with st.spinner(
            "Searching documents and generating answer..."
        ):

            try:

                (
                    answer,
                    sources,
                    source_type,
                ) = perform_rag(prompt)

                st.markdown(answer)

                # --------------------------------------------
                # Evidence
                # --------------------------------------------

                if sources:

                    with st.expander(
                        "🔎 Retrieved evidence"
                    ):

                        st.caption(
                            f"Source type: {source_type}"
                        )

                        for index, source in enumerate(
                            sources,
                            start=1,
                        ):

                            if not isinstance(
                                source,
                                dict,
                            ):
                                continue

                            filename = (
                                source.get(
                                    "filename"
                                )
                                or source.get(
                                    "source"
                                )
                                or "Unknown"
                            )

                            page = source.get(
                                "page"
                            )

                            similarity = (
                                source.get(
                                    "similarity_score"
                                )
                            )

                            st.markdown(
                                f"**{index}. {filename}**"
                            )

                            if page is not None:

                                st.caption(
                                    f"Page/Sheet: {page}"
                                )

                            if similarity is not None:

                                try:

                                    st.caption(
                                        "Relevance: "
                                        f"{float(similarity):.3f}"
                                    )

                                except Exception:
                                    pass

                            url = source.get(
                                "url"
                            )

                            if url:

                                st.caption(
                                    f"URL: {url}"
                                )

                            content = (
                                source.get(
                                    "content"
                                )
                                or source.get(
                                    "text"
                                )
                                or ""
                            )

                            if content:

                                st.text(
                                    str(content)[:2000]
                                )

                            st.divider()

                # --------------------------------------------
                # Save assistant message.
                # --------------------------------------------

                st.session_state.db.save_message(
                    conversation_id,
                    "assistant",
                    answer,
                    user_id=user_id,
                )

                # --------------------------------------------
                # Auto-title conversation.
                # --------------------------------------------

                current_messages = (
                    st.session_state.db.get_messages(
                        conversation_id,
                        user_id=user_id,
                    )
                )

                if len(current_messages) <= 2:

                    title = prompt[:70]

                    if len(prompt) > 70:
                        title += "..."

                    st.session_state.db.update_conversation_title(
                        conversation_id,
                        title,
                        user_id=user_id,
                    )

            except Exception as exc:

                answer = (
                    "Something went wrong while "
                    "processing your request."
                )

                st.error(answer)
                st.exception(exc)

                try:

                    st.session_state.db.save_message(
                        conversation_id,
                        "assistant",
                        f"{answer}\n\n{exc}",
                        user_id=user_id,
                    )

                except Exception:
                    pass


# ============================================================
# DOCUMENT PAGE
# ============================================================

def documents_page():

    st.title(
        "📚 Document Management"
    )

    user = st.session_state.user

    is_admin = (
        str(
            user["role"]
        ).lower()
        == "admin"
    )

    # ========================================================
    # UPLOAD
    # ========================================================

    st.subheader(
        "Upload documents"
    )

    uploaded_files = st.file_uploader(
        "Choose one or more files",
        type=[
            item.lstrip(".")
            for item in sorted(
                SUPPORTED_EXTENSIONS
            )
        ],
        accept_multiple_files=True,
    )

    col1, col2, col3 = st.columns(3)

    with col1:

        chunk_method = st.selectbox(
            "Chunking method",
            [
                "Character",
                "Recursive",
                "Token",
                "Markdown",
                "Context",
                "Multimodal",
            ],
            index=1,
        )

    with col2:

        chunk_size = st.number_input(
            "Chunk size",
            min_value=100,
            max_value=10000,
            value=int(
                DEFAULT_CHUNK_SIZE
            ),
            step=100,
        )

    with col3:

        max_overlap = max(
            0,
            int(chunk_size) - 1,
        )

        default_overlap = min(
            int(DEFAULT_CHUNK_OVERLAP),
            max_overlap,
        )

        chunk_overlap = st.number_input(
            "Chunk overlap",
            min_value=0,
            max_value=max_overlap,
            value=default_overlap,
            step=(
                50
                if max_overlap >= 50
                else 1
            ),
        )

    # ========================================================
    # PROCESS UPLOAD
    # ========================================================

    if uploaded_files:

        if st.button(
            "🚀 Process and upload",
            type="primary",
            use_container_width=True,
        ):

            # ------------------------------------------------
            # Refresh user.
            # ------------------------------------------------

            fresh_user = (
                st.session_state.db.get_user_by_id(
                    user["id"]
                )
            )

            if not fresh_user:

                st.error(
                    "Your account no longer exists. "
                    "Please log in again."
                )

                st.session_state.authenticated = False
                st.session_state.user = None
                st.session_state.conversation_id = None
                st.session_state.selected_document_ids = []

                st.rerun()

            st.session_state.user = fresh_user

            user = fresh_user

            upload_success = False

            # ------------------------------------------------
            # Process each file.
            # ------------------------------------------------

            for uploaded_file in uploaded_files:

                filename = (
                    uploaded_file.name
                )

                try:

                    file_bytes = (
                        uploaded_file.getvalue()
                    )

                    extension = Path(
                        filename
                    ).suffix.lower()

                    if extension not in SUPPORTED_EXTENSIONS:

                        st.error(
                            f"{filename}: unsupported file type."
                        )

                        continue

                    with st.status(
                        f"Processing {filename}...",
                        expanded=True,
                    ):

                        # ------------------------------------
                        # Extraction
                        # ------------------------------------

                        st.write(
                            "1️⃣ Extracting document..."
                        )

                        extracted = extract_document(
                            filename,
                            file_bytes,
                        )

                        if not extracted:

                            raise ValueError(
                                "No text or data could be extracted."
                            )

                        st.write(
                            f"Extracted {len(extracted)} section(s)."
                        )

                        # ------------------------------------
                        # Chunking
                        # ------------------------------------

                        st.write(
                            "2️⃣ Creating chunks..."
                        )

                        chunks = chunk_document(
                            extracted,
                            chunk_method,
                            int(chunk_size),
                            int(chunk_overlap),
                            filename,
                        )

                        if not chunks:

                            raise ValueError(
                                "Chunking produced zero usable chunks."
                            )

                        st.write(
                            f"Created {len(chunks)} chunk(s)."
                        )

                        # ------------------------------------
                        # Save
                        # ------------------------------------

                        st.write(
                            "3️⃣ Saving document and chunks..."
                        )

                        result = (
                            st.session_state.db
                            .save_document_with_chunks(
                                user_id=int(user["id"]),
                                filename=filename,
                                file_type=extension,
                                file_bytes=file_bytes,
                                chunking_method=chunk_method,
                                chunk_size=int(chunk_size),
                                chunk_overlap=int(chunk_overlap),
                                chunks=chunks,
                                metadata={
                                    "original_filename": filename,
                                    "sections": len(extracted),
                                    "chunks": len(chunks),
                                    "chunking_method": chunk_method,
                                },
                            )
                        )

                        st.write(
                            "4️⃣ PostgreSQL transaction committed."
                        )

                        st.write(
                            f"Document ID: "
                            f"{result['document_id']}"
                        )

                        st.write(
                            f"Chunks saved: "
                            f"{result['chunk_count']}"
                        )

                    st.success(
                        f"✅ {filename} uploaded and indexed successfully."
                    )

                    upload_success = True

                except Exception as exc:

                    st.error(
                        f"❌ Failed to process "
                        f"{filename}: {exc}"
                    )

                    with st.expander(
                        "Technical error"
                    ):

                        st.exception(exc)

            if upload_success:

                st.session_state.selected_document_ids = []
                st.session_state.document_search = ""

                st.rerun()

    st.divider()

    # ========================================================
    # DOCUMENT SEARCH
    # ========================================================

    st.subheader(
        "🔎 Find a document"
    )

    search_col, clear_col = st.columns(
        [5, 1]
    )

    with search_col:

        document_search = st.text_input(
            "Search by filename",
            value=(
                st.session_state.document_search
            ),
            placeholder=(
                "Example: report.pdf..."
            ),
            label_visibility="collapsed",
        )

        st.session_state.document_search = (
            document_search
        )

    with clear_col:

        if st.button(
            "Clear",
            use_container_width=True,
        ):

            st.session_state.document_search = ""

            st.rerun()

    # ========================================================
    # GET DOCUMENTS
    # ========================================================

    try:

        documents = (
            st.session_state.db.get_documents(
                user_id=int(user["id"]),
                is_admin=is_admin,
                search=document_search,
            )
        )

    except Exception as exc:

        st.error(
            "Could not load documents."
        )

        st.exception(exc)

        return

    if is_admin:

        st.subheader(
            f"All documents ({len(documents)})"
        )

    else:

        st.subheader(
            f"Your documents ({len(documents)})"
        )

    if not documents:

        if document_search:

            st.info(
                f'No documents match "{document_search}".'
            )

        else:

            st.info(
                "No documents have been uploaded yet."
            )

        return

    # ========================================================
    # RETRIEVAL DOCUMENT SELECTION
    # ========================================================

    st.subheader(
        "🎯 Retrieval document selection"
    )

    document_options = {
        f"{doc['filename']}  (#{doc['id']})":
        doc["id"]
        for doc in documents
    }

    valid_ids = set(
        document_options.values()
    )

    cleaned_selected_ids = []

    for document_id in (
        st.session_state.selected_document_ids
    ):

        try:

            document_id = int(
                document_id
            )

        except Exception:

            continue

        if document_id in valid_ids:
            cleaned_selected_ids.append(
                document_id
            )

    st.session_state.selected_document_ids = (
        cleaned_selected_ids
    )

    selected_labels_default = [
        label
        for label, document_id
        in document_options.items()
        if document_id
        in st.session_state.selected_document_ids
    ]

    selected_labels = st.multiselect(
        "Choose documents that Chat should search",
        options=list(
            document_options.keys()
        ),
        default=selected_labels_default,
        help=(
            "Select one or more documents. "
            "Chat will search only those documents. "
            "Leave empty to search all documents."
        ),
    )

    st.session_state.selected_document_ids = [
        document_options[label]
        for label in selected_labels
    ]

    if selected_labels:

        st.success(
            "🎯 Chat is restricted to: "
            + ", ".join(
                selected_labels
            )
        )

    else:

        st.caption(
            "No document filter is active. "
            "Chat searches across all your documents."
        )

    # ========================================================
    # CHUNK TYPE FILTER
    # ========================================================

    chunk_type_options = [
        "Recursive",
        "Character",
        "Token",
        "Markdown",
        "Context",
        "Table",
        "Image",
        "Visual",
    ]

    selected_chunk_types = st.multiselect(
        "Optional chunk-type filter",
        chunk_type_options,
        default=(
            st.session_state.selected_chunk_types
        ),
        help=(
            "Leave empty to search all chunk types."
        ),
    )

    st.session_state.selected_chunk_types = (
        selected_chunk_types
    )

    # ========================================================
    # DOCUMENT LIST
    # ========================================================

    st.subheader(
        "📄 Uploaded documents"
    )

    for document in documents:

        document_id = document["id"]
        filename = document["filename"]

        size = document.get(
            "file_size",
            0,
        )

        try:

            size_mb = (
                float(size or 0)
                / 1024
                / 1024
            )

        except Exception:

            size_mb = 0.0

        chunk_count = document.get(
            "chunk_count",
            0,
        )

        is_selected = (
            document_id
            in st.session_state.selected_document_ids
        )

        with st.container(
            border=True
        ):

            col1, col2, col3 = st.columns(
                [6, 2, 1]
            )

            with col1:

                prefix = (
                    "🎯 📄"
                    if is_selected
                    else "📄"
                )

                st.markdown(
                    f"### {prefix} {filename}"
                )

                if is_admin:

                    st.caption(
                        "Owner: "
                        f"{document.get('username', 'Unknown')}"
                    )

                st.caption(
                    f"ID: #{document_id} "
                    f"• Type: {document.get('file_type', '')} "
                    f"• Size: {size_mb:.2f} MB"
                )

                st.caption(
                    f"Chunks: {chunk_count} "
                    f"• Chunking: "
                    f"{document.get('chunking_method', '')}"
                )

            with col2:

                st.caption(
                    "Created:"
                )

                st.caption(
                    str(
                        document.get(
                            "created_at",
                            "",
                        )
                    )
                )

            with col3:

                if st.button(
                    "🗑️ Delete",
                    key=(
                        f"delete_doc_"
                        f"{document_id}"
                    ),
                ):

                    try:

                        deleted = (
                            st.session_state.db
                            .delete_document(
                                document_id,
                                user_id=user["id"],
                                is_admin=is_admin,
                            )
                        )

                        if deleted:

                            if (
                                document_id
                                in st.session_state.selected_document_ids
                            ):

                                st.session_state.selected_document_ids.remove(
                                    document_id
                                )

                            st.success(
                                "Document deleted."
                            )

                            st.rerun()

                        else:

                            st.error(
                                "Could not delete document."
                            )

                    except Exception as exc:

                        st.error(
                            f"Delete failed: {exc}"
                        )


# ============================================================
# ADMIN PAGE
# ============================================================

def admin_page():

    user = st.session_state.user

    if str(
        user["role"]
    ).lower() != "admin":

        st.error(
            "You do not have administrator access."
        )

        return

    st.title(
        "🛠️ Administration"
    )

    # ========================================================
    # STATS
    # ========================================================

    try:

        stats = (
            st.session_state.db.get_stats()
        )

    except Exception as exc:

        st.error(
            "Could not load statistics."
        )

        st.exception(exc)

        stats = {}

    columns = st.columns(5)

    for column, (
        label,
        key,
    ) in zip(
        columns,
        [
            ("Users", "users"),
            ("Documents", "documents"),
            ("Chunks", "chunks"),
            ("Chats", "conversations"),
            ("Messages", "messages"),
        ],
    ):

        with column:

            st.metric(
                label,
                stats.get(
                    key,
                    0,
                ),
            )

    st.divider()

    # ========================================================
    # USER MANAGEMENT
    # ========================================================

    st.subheader(
        "User management"
    )

    try:

        users = (
            st.session_state.db.get_users()
        )

    except Exception as exc:

        st.error(
            "Could not load users."
        )

        st.exception(exc)

        return

    for target in users:

        target_id = target["id"]

        target_role = str(
            target.get(
                "role",
                "user",
            )
        ).lower()

        with st.container(
            border=True
        ):

            col1, col2, col3, col4 = (
                st.columns(
                    [3, 2, 2, 2]
                )
            )

            with col1:

                st.markdown(
                    f"**{target['username']}**"
                )

                st.caption(
                    f"User ID: {target_id}"
                )

            with col2:

                st.write(
                    f"Role: `{target_role}`"
                )

            with col3:

                new_role = st.selectbox(
                    "Role",
                    [
                        "user",
                        "admin",
                    ],
                    index=(
                        1
                        if target_role == "admin"
                        else 0
                    ),
                    key=(
                        f"role_"
                        f"{target_id}"
                    ),
                )

                if (
                    new_role
                    != target_role
                ):

                    if st.button(
                        "Update role",
                        key=(
                            f"update_role_"
                            f"{target_id}"
                        ),
                    ):

                        try:

                            changed = (
                                st.session_state.db
                                .set_user_role(
                                    target_id,
                                    new_role,
                                )
                            )

                            if changed:

                                if (
                                    target_id
                                    == user["id"]
                                ):

                                    fresh = (
                                        st.session_state.db
                                        .get_user_by_id(
                                            user["id"]
                                        )
                                    )

                                    if fresh:

                                        st.session_state.user = (
                                            fresh
                                        )

                                st.success(
                                    "Role updated."
                                )

                                st.rerun()

                            else:

                                st.error(
                                    "Cannot remove the last administrator."
                                )

                        except Exception as exc:

                            st.error(
                                str(exc)
                            )

            with col4:

                with st.expander(
                    "Password"
                ):

                    new_password = st.text_input(
                        "New password",
                        type="password",
                        key=(
                            f"password_"
                            f"{target_id}"
                        ),
                    )

                    if st.button(
                        "Reset password",
                        key=(
                            f"reset_"
                            f"{target_id}"
                        ),
                    ):

                        if len(
                            new_password
                        ) < 6:

                            st.error(
                                "Password must be at least 6 characters."
                            )

                        else:

                            try:

                                (
                                    st.session_state.db
                                    .reset_user_password(
                                        target_id,
                                        new_password,
                                    )
                                )

                                st.success(
                                    "Password reset."
                                )

                            except Exception as exc:

                                st.error(
                                    str(exc)
                                )

                if target_id != user["id"]:

                    if st.button(
                        "Delete user",
                        key=(
                            f"delete_user_"
                            f"{target_id}"
                        ),
                    ):

                        try:

                            deleted = (
                                st.session_state.db
                                .delete_user(
                                    target_id
                                )
                            )

                            if deleted:

                                st.success(
                                    "User deleted."
                                )

                                st.rerun()

                            else:

                                st.error(
                                    "Could not delete user. "
                                    "The last administrator cannot be deleted."
                                )

                        except Exception as exc:

                            st.error(
                                str(exc)
                            )


# ============================================================
# MAIN APPLICATION
# ============================================================

if not st.session_state.authenticated:

    login_page()

    st.stop()


# ============================================================
# VALIDATE SESSION AGAINST DATABASE
# ============================================================

if not validate_logged_in_user():

    st.warning(
        "Your login session is no longer valid. "
        "Please log in again."
    )

    st.session_state.authenticated = False
    st.session_state.user = None
    st.session_state.conversation_id = None
    st.session_state.selected_document_ids = []

    st.stop()


# ============================================================
# INITIALIZE LLM
# ============================================================

if st.session_state.llm is None:

    st.session_state.llm = create_llm()


# ============================================================
# ROUTING
# ============================================================

page = render_sidebar()

if page == "💬 Chat":

    chat_page()

elif page == "📚 Documents":

    documents_page()

elif page == "🛠️ Admin":

    admin_page()